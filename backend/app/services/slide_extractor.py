"""Extract text from optional PPT/PDF supporting material."""
import os


def extract_slides_text(path: str) -> str:
    if not path or not os.path.exists(path):
        return ""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for i, slide in enumerate(prs.slides, 1):
                parts = [sh.text.strip() for sh in slide.shapes
                         if hasattr(sh, "text") and sh.text.strip()]
                if parts:
                    texts.append(f"Slide {i}: " + " | ".join(parts))
            return "\n".join(texts)[:4000]
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(path)
            texts = []
            for i, page in enumerate(reader.pages[:15], 1):
                t = (page.extract_text() or "").strip()
                if t:
                    texts.append(f"Page {i}: {t[:400]}")
            return "\n".join(texts)[:4000]
    except Exception as e:
        print(f"[SlideExtractor] Failed to read slides: {e}")
    return ""
